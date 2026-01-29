import os
import sys
import zipfile
import shutil
import re
import warnings
import json
from pathlib import Path
from pptx import Presentation
import whisper
from faster_whisper import WhisperModel
import torch
from tqdm import tqdm

# Fix Windows console encoding issues
if sys.platform == 'win32':
    sys.stdout.reconfigure(encoding='utf-8')

# Suppress CUDA/Triton warnings for cleaner output
warnings.filterwarnings("ignore", category=UserWarning, module="whisper")
warnings.filterwarnings("ignore", category=UserWarning, module="faster_whisper")

# ⚙️ CONFIGURATION SETTINGS - Edit these for easy customization
# 📂 Folder Settings
PPTX_FOLDER = "presentations"   # input folder
OUTPUT_FOLDER = "output"        # output folder

# 🚀 Transcription Engine Selection
# "standard" = openai-whisper (original implementation, reliable)
# "faster-whisper" = 4-5x faster, uses less memory, supports INT8 quantization on CPU
TRANSCRIPTION_ENGINE = "faster-whisper"  # Options: "standard", "faster-whisper" (recommended)

# 🎯 Whisper Model Settings
WHISPER_MODEL = "base"        # Options: "tiny", "base", "small", "medium", "large"
FORCE_LANGUAGE = "en"           # Force language to prevent mixing (None for auto-detect)

# ⚡ Performance Settings
FORCE_DEVICE = "cpu"             # Options: None (auto), "cpu", "cuda" (force specific device)
USE_HALF_PRECISION = False       # fp16 for 30-50% speed boost (minimal accuracy loss)
GPU_BEST_OF = 3                 # Decoding attempts on GPU (higher = more accurate, slower)
GPU_BEAM_SIZE = 3               # Beam search size on GPU
CPU_BEST_OF = 3                 # Decoding attempts on CPU
CPU_BEAM_SIZE = 3               # Beam search size on CPU

# 🎚️ Quality Settings
TEMPERATURE = 0.0               # 0.0 = deterministic, 0.1-1.0 = more creative
ENABLE_WORD_TIMESTAMPS = True   # Get word-level timing data

os.makedirs(OUTPUT_FOLDER, exist_ok=True)

# ⚡ Load Whisper model with optimal device selection
def get_optimal_device():
    # Check if user forced a specific device
    if FORCE_DEVICE:
        if FORCE_DEVICE == "cuda" and not torch.cuda.is_available():
            print("[!] CUDA requested but not available, falling back to CPU")
            return "cpu"
        print(f"[*] Forced device: {FORCE_DEVICE}")
        return FORCE_DEVICE

    # Auto-detect best device
    if torch.cuda.is_available():
        gpu_name = torch.cuda.get_device_name(0)
        gpu_memory = torch.cuda.get_device_properties(0).total_memory / 1024**3
        print(f"GPU detected: {gpu_name} ({gpu_memory:.1f}GB)")
        return "cuda"
    else:
        print("No GPU detected, using CPU")
        return "cpu"

device = get_optimal_device()

# ⚡ Load the appropriate Whisper model based on selected engine
model = None
faster_model = None

if TRANSCRIPTION_ENGINE == "faster-whisper":
    print(f"Loading faster-whisper model ({WHISPER_MODEL}) on {device}...")
    # For faster-whisper, we need to specify compute type
    compute_type = "float16" if device == "cuda" and USE_HALF_PRECISION else "int8" if device == "cpu" else "float16"
    faster_model = WhisperModel(WHISPER_MODEL, device=device, compute_type=compute_type)
    print(f"[OK] Using faster-whisper with {compute_type} precision")
elif TRANSCRIPTION_ENGINE == "standard":
    print(f"Loading standard openai-whisper model ({WHISPER_MODEL}) on {device}...")
    model = whisper.load_model(WHISPER_MODEL, device=device)
    print(f"[OK] Using standard openai-whisper")
else:
    print(f"[!] Invalid TRANSCRIPTION_ENGINE '{TRANSCRIPTION_ENGINE}', falling back to standard")
    TRANSCRIPTION_ENGINE = "standard"
    model = whisper.load_model(WHISPER_MODEL, device=device)

def extract_text_from_pptx(pptx_path):
    """Extract all text from slides in a PPTX."""
    prs = Presentation(pptx_path)
    texts = []
    for slide_num, slide in enumerate(prs.slides, start=1):
        slide_texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text.strip()
                if text:
                    slide_texts.append(text)
        if slide_texts:
            texts.append(f"--- Slide {slide_num} ---\n" + "\n".join(slide_texts))
    return "\n\n".join(texts)

def extract_audio_from_pptx(pptx_path, temp_dir):
    """Extract embedded audio files from PPTX (wav, mp3, m4a) with proper ordering."""
    audio_files = []
    with zipfile.ZipFile(pptx_path, "r") as zip_ref:
        media_files = []
        for file in zip_ref.namelist():
            if file.startswith("ppt/media/") and file.lower().endswith((".wav", ".mp3", ".m4a")):
                media_files.append(file)

        # Sort by the numeric part in filename (media1, media2, etc.)
        def get_media_number(filename):
            match = re.search(r'media(\d+)', filename)
            return int(match.group(1)) if match else 0

        media_files.sort(key=get_media_number)

        for file in media_files:
            extracted_path = os.path.join(temp_dir, os.path.basename(file))
            with open(extracted_path, "wb") as f:
                f.write(zip_ref.read(file))
            audio_files.append(extracted_path)

    return audio_files


def get_checkpoint_file(audio_path):
    """Get checkpoint file path for a specific audio file."""
    base_name = Path(audio_path).stem
    checkpoint_file = Path(OUTPUT_FOLDER) / f"{base_name}_checkpoint.json"
    return checkpoint_file

def load_checkpoint(checkpoint_file):
    """Load existing checkpoint if available."""
    if checkpoint_file.exists():
        with open(checkpoint_file, 'r', encoding='utf-8') as f:
            return json.load(f)
    return None

def save_checkpoint(checkpoint_file, checkpoint_data):
    """Save checkpoint data to file."""
    with open(checkpoint_file, 'w', encoding='utf-8') as f:
        json.dump(checkpoint_data, f, ensure_ascii=False, indent=2)

def transcribe_single_file(audio_path):
    """Transcribe a single audio/video file using the selected engine with checkpoint support."""
    checkpoint_file = get_checkpoint_file(audio_path)
    
    if TRANSCRIPTION_ENGINE == "faster-whisper":
        # Check for existing checkpoint
        checkpoint = load_checkpoint(checkpoint_file)
        if checkpoint:
            last_segment_idx = len(checkpoint.get("segments", []))
            last_timestamp = checkpoint["segments"][-1]["end"] if checkpoint.get("segments") else 0
            print(f"[Checkpoint] Found existing progress! {last_segment_idx} segments ({last_timestamp:.1f}s)")
        else:
            checkpoint = {"segments": [], "metadata": {}}
            last_segment_idx = 0
            last_timestamp = 0
            print(f"[Transcribing] Starting new transcription...")
        
        segments, info = faster_model.transcribe(
            audio_path,
            language=FORCE_LANGUAGE,
            task="transcribe",
            temperature=TEMPERATURE,
            beam_size=GPU_BEAM_SIZE if device == "cuda" else CPU_BEAM_SIZE,
            word_timestamps=ENABLE_WORD_TIMESTAMPS
        )
        
        # Store metadata
        checkpoint["metadata"] = {
            "duration": info.duration,
            "language": info.language,
            "file": str(audio_path)
        }
        
        print(f"[Info] Duration: {info.duration:.1f}s | Language: {info.language}")
        if last_segment_idx == 0:
            print(f"[Checkpoint] Saving to: {checkpoint_file}")
        
        # Create progress bar
        with tqdm(total=int(info.duration), desc="Transcribing", unit="s", 
                  bar_format="{l_bar}{bar}| {n:.0f}/{total:.0f}s [{elapsed}<{remaining}]",
                  initial=int(last_timestamp)) as pbar:
            
            segment_index = 0
            last_position = int(last_timestamp)
            
            for segment in segments:
                # Skip already processed segments
                if segment_index < last_segment_idx:
                    segment_index += 1
                    continue
                
                # Add new segment to checkpoint
                segment_data = {
                    'start': segment.start,
                    'end': segment.end,
                    'text': segment.text
                }
                checkpoint["segments"].append(segment_data)
                
                # Save checkpoint every 10 segments (balance between safety and I/O)
                if segment_index % 10 == 0 or segment_index == last_segment_idx:
                    save_checkpoint(checkpoint_file, checkpoint)
                
                # Update progress bar
                current_position = int(segment.end)
                pbar.update(current_position - last_position)
                last_position = current_position
                segment_index += 1
            
            # Final save
            save_checkpoint(checkpoint_file, checkpoint)
        
        print(f"[Completed] Transcription finished! Total segments: {len(checkpoint['segments'])}")
        
        # Extract text from checkpoint
        text = " ".join([seg["text"] for seg in checkpoint["segments"]])
        
        # Clean up checkpoint on successful completion
        print(f"[Cleanup] Removing checkpoint file...")
        checkpoint_file.unlink(missing_ok=True)
        
        return {"text": text.strip()}
    else:
        # Standard openai-whisper
        if device == "cuda":
            result = model.transcribe(
                audio_path,
                language=FORCE_LANGUAGE,
                task="transcribe",
                temperature=TEMPERATURE,
                best_of=GPU_BEST_OF,
                beam_size=GPU_BEAM_SIZE,
                word_timestamps=ENABLE_WORD_TIMESTAMPS,
                fp16=USE_HALF_PRECISION
            )
        else:
            result = model.transcribe(
                audio_path,
                language=FORCE_LANGUAGE,
                task="transcribe",
                temperature=TEMPERATURE,
                best_of=CPU_BEST_OF,
                beam_size=CPU_BEAM_SIZE,
                word_timestamps=ENABLE_WORD_TIMESTAMPS,
                fp16=False
            )
        return result

def transcribe_audio(audio_files):
    """Run Whisper transcription with hybrid CPU/GPU optimization and progress bar."""
    transcripts = []

    # Create progress bar
    progress_bar = tqdm(
        audio_files,
        desc="[Audio] Transcribing",
        unit="file",
        bar_format="{l_bar}{bar}| {n_fmt}/{total_fmt} [{elapsed}<{remaining}]"
    )

    for audio_path in progress_bar:
        # Update progress bar description with current file
        filename = os.path.basename(audio_path)
        media_match = re.search(r'media(\d+)', filename)
        media_num = media_match.group(1) if media_match else "unknown"

        progress_bar.set_description(f"[Audio] Processing {media_num}")

        try:
            # Use the unified transcribe function
            result = transcribe_single_file(audio_path)

        except Exception as e:
            progress_bar.write(f"[!] Error with {filename}: {e}")
            # Fallback - try with standard whisper if faster-whisper fails
            if TRANSCRIPTION_ENGINE == "faster-whisper":
                progress_bar.write(f"[!] Falling back to standard whisper for {filename}...")
                cpu_model = whisper.load_model(WHISPER_MODEL, device="cpu")
                result = cpu_model.transcribe(
                    audio_path,
                    language=FORCE_LANGUAGE,
                    temperature=TEMPERATURE
                )
            else:
                raise e

        transcripts.append(f"--- Audio {media_num} Transcript ---\n{result['text'].strip()}")

        # Clear GPU cache periodically if using CUDA
        if device == "cuda" and torch.cuda.is_available():
            torch.cuda.empty_cache()

    progress_bar.close()
    return "\n\n".join(transcripts)

def process_pptx(pptx_path):
    """Process one PowerPoint: text + audio → output TXT file."""
    print(f"\n[PPTX] Processing {pptx_path}...")
    base_name = os.path.splitext(os.path.basename(pptx_path))[0]
    temp_dir = os.path.join(OUTPUT_FOLDER, base_name + "_media")
    os.makedirs(temp_dir, exist_ok=True)

    try:
        # Slide text
        text_content = extract_text_from_pptx(pptx_path)

        # Embedded audio
        audio_files = extract_audio_from_pptx(pptx_path, temp_dir)
        transcript = transcribe_audio(audio_files) if audio_files else ""

        # Combine output with improved organization
        final_output = []
        if text_content:
            final_output.append("### PowerPoint Slide Content (In Order) ###\n" + text_content)
        if transcript:
            final_output.append("\n### Audio Transcripts (In Chronological Order) ###\n" + transcript)

        # Add summary note about ordering
        if text_content and transcript:
            final_output.append("\n### Note ###\nAudio transcripts are generated by OpenAI Whisper.")

        # Save result
        output_path = os.path.join(OUTPUT_FOLDER, base_name + ".txt")
        with open(output_path, "w", encoding="utf-8") as f:
            f.write("\n\n".join(final_output))

        print(f"[OK] Saved results to {output_path}")

    finally:
        # Cleanup extracted media
        shutil.rmtree(temp_dir, ignore_errors=True)

def process_mp3(mp3_path):
    """Process standalone MP3 file: transcribe audio → output TXT file."""
    print(f"\n[MP3] Processing {mp3_path}...")
    base_name = os.path.splitext(os.path.basename(mp3_path))[0]

    try:
        # Transcribe the MP3 file
        print("[Audio] Transcribing...")
        result = transcribe_single_file(mp3_path)

        # Save transcription
        output_path = os.path.join(OUTPUT_FOLDER, base_name + ".txt")
        with open(output_path, "w", encoding="utf-8") as f:
            f.write(f"### MP3 Audio Transcription ###\n\n{result['text'].strip()}")

        print(f"[OK] Saved transcription to {output_path}")

        # Clear GPU cache if using CUDA
        if device == "cuda" and torch.cuda.is_available():
            torch.cuda.empty_cache()

    except Exception as e:
        print(f"[ERROR] Error processing {mp3_path}: {e}")

def process_mp4(mp4_path):
    """Process MP4 video file: extract and transcribe audio → output TXT file."""
    print(f"\n[MP4] Processing {mp4_path}...")
    base_name = os.path.splitext(os.path.basename(mp4_path))[0]

    try:
        # Transcribe the MP4 file (Whisper extracts audio automatically)
        print("[Audio] Extracting and transcribing...")
        result = transcribe_single_file(mp4_path)

        # Save transcription
        output_path = os.path.join(OUTPUT_FOLDER, base_name + ".txt")
        with open(output_path, "w", encoding="utf-8") as f:
            f.write(result['text'].strip())

        print(f"[OK] Saved transcription to {output_path}")

        # Clear GPU cache if using CUDA
        if device == "cuda" and torch.cuda.is_available():
            torch.cuda.empty_cache()

    except Exception as e:
        print(f"[ERROR] Error processing {mp4_path}: {e}")

def main():
    # Scan for all supported file types
    pptx_files = [f for f in os.listdir(PPTX_FOLDER) if f.lower().endswith(".pptx")]
    mp3_files = [f for f in os.listdir(PPTX_FOLDER) if f.lower().endswith(".mp3")]
    mp4_files = [f for f in os.listdir(PPTX_FOLDER) if f.lower().endswith(".mp4")]

    total_files = len(pptx_files) + len(mp3_files) + len(mp4_files)

    if total_files == 0:
        print(f"[!] No .pptx, .mp3, or .mp4 files found in {PPTX_FOLDER}")
        return

    print(f"[Files] Found {len(pptx_files)} PPTX, {len(mp3_files)} MP3, {len(mp4_files)} MP4 files")

    # Process PPTX files
    for file in pptx_files:
        process_pptx(os.path.join(PPTX_FOLDER, file))

    # Process MP3 files
    for file in mp3_files:
        process_mp3(os.path.join(PPTX_FOLDER, file))

    # Process MP4 files
    for file in mp4_files:
        process_mp4(os.path.join(PPTX_FOLDER, file))

if __name__ == "__main__":
    main()
